"""
Script para generar la presentación "MLMonitor en CML: Por qué ir 100% Cloud".
Ejecutar con: poetry run python scripts/create_cml_presentation.py
Output: artifacts/ppt/MLMonitor_CML_Decision.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paleta (misma que create_presentation.py)
# ---------------------------------------------------------------------------
BG_DARK      = RGBColor(0x0F, 0x0F, 0x0E)
BG_CARD      = RGBColor(0x1C, 0x1C, 0x1A)
ACCENT       = RGBColor(0xD9, 0x77, 0x06)
ACCENT_LIGHT = RGBColor(0xF5, 0xA6, 0x23)
TEXT_CREAM   = RGBColor(0xE8, 0xE3, 0xDB)
TEXT_GRAY    = RGBColor(0xB0, 0xAD, 0xA8)
TEXT_DIM     = RGBColor(0x6B, 0x68, 0x65)
RED          = RGBColor(0xEF, 0x44, 0x44)
RED_DARK     = RGBColor(0x2A, 0x10, 0x10)
RED_SOFT     = RGBColor(0xF5, 0xA0, 0xA0)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
YELLOW       = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers base
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide


def _txt(slide, text, left, top, width, height,
         size=14, bold=False, color=TEXT_CREAM,
         align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box


def _rect(slide, left, top, width, height, fill_color,
          line_color=None, line_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _badge(slide, label):
    _rect(slide, Inches(0.5), Inches(0.28), Inches(4.0), Inches(0.32),
          RGBColor(0x2A, 0x1E, 0x08), line_color=ACCENT, line_pt=0.5)
    _txt(slide, label,
         Inches(0.6), Inches(0.28), Inches(3.9), Inches(0.32),
         size=7.5, bold=True, color=ACCENT)


def _divider(slide, y=Inches(1.88)):
    _rect(slide, Inches(0.5), y, SLIDE_W - Inches(1.0), Inches(0.018), TEXT_DIM)


def _page_num(slide, n):
    _txt(slide, str(n),
         SLIDE_W - Inches(0.6), SLIDE_H - Inches(0.38),
         Inches(0.5), Inches(0.3),
         size=9, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def _bullets(slide, items, left, top, width, spacing=0.82, size=12, accent_color=None):
    """
    Renders a list of (bold_keyword, description) as styled bullets.
    Style: ▸  Bold keyword  —  descripción en gris
    """
    dot_color = accent_color or ACCENT
    for i, (bold_text, desc_text) in enumerate(items):
        y = top + Inches(i * spacing)
        box = slide.shapes.add_textbox(left, y, width, Inches(spacing * 0.92))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        r0 = p.add_run()
        r0.text = "▸  "
        r0.font.size = Pt(size)
        r0.font.bold = False
        r0.font.color.rgb = dot_color
        r0.font.name = "Calibri"

        r1 = p.add_run()
        r1.text = bold_text
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = TEXT_CREAM
        r1.font.name = "Calibri"

        r2 = p.add_run()
        r2.text = "  —  "
        r2.font.size = Pt(size)
        r2.font.bold = False
        r2.font.color.rgb = TEXT_DIM
        r2.font.name = "Calibri"

        r3 = p.add_run()
        r3.text = desc_text
        r3.font.size = Pt(size)
        r3.font.bold = False
        r3.font.color.rgb = TEXT_GRAY
        r3.font.name = "Calibri"


# ---------------------------------------------------------------------------
# Slide 1 — Portada
# ---------------------------------------------------------------------------

def slide_portada(prs):
    s = _blank_slide(prs)

    # Barra izquierda naranja
    _rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)

    _txt(s, "MLMonitor en CML",
         Inches(1.2), Inches(1.5), Inches(11.5), Inches(1.5),
         size=50, bold=True, color=TEXT_CREAM)

    _txt(s, "Por qué una arquitectura híbrida suma más problemas de los que resuelve",
         Inches(1.2), Inches(3.05), Inches(10.8), Inches(0.85),
         size=18, color=ACCENT)

    _rect(s, Inches(1.2), Inches(3.88), Inches(5.0), Inches(0.02), TEXT_DIM)

    _txt(s, "Equipo de Crédito  ·  BazBoost V1  ·  Abril 2026",
         Inches(1.2), Inches(4.02), Inches(11), Inches(0.45),
         size=12, color=TEXT_DIM)

    # Pills de los tres temas de la presentación
    topics = ["Dependencia operativa", "Complejidad técnica", "Trabajo transitorio"]
    for i, topic in enumerate(topics):
        x = Inches(1.2 + i * 3.85)
        _rect(s, x, Inches(5.1), Inches(3.6), Inches(0.48),
              BG_CARD, line_color=ACCENT, line_pt=0.75)
        _txt(s, topic, x, Inches(5.1), Inches(3.6), Inches(0.48),
             size=11, bold=True, color=TEXT_CREAM, align=PP_ALIGN.CENTER)

    _page_num(s, 1)


# ---------------------------------------------------------------------------
# Slide 2 — Contexto: arquitectura híbrida propuesta
# ---------------------------------------------------------------------------

def slide_arquitectura(prs):
    s = _blank_slide(prs)
    _badge(s, "CONTEXTO · ARQUITECTURA PROPUESTA")

    _txt(s, "El flujo híbrido CML + Nube",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    _txt(s, "ETL on-premise en CML  ·  Pipeline en AWS  ·  sincronización semanal manual",
         Inches(0.5), Inches(1.42), Inches(12), Inches(0.4),
         size=12, color=TEXT_DIM)

    _divider(s)

    # Diagrama de flujo horizontal
    nodes = [
        ("CML · VM",    "ETL semanal\nCSV → PostgreSQL local",      ACCENT_LIGHT),
        ("SYNC\nMANUAL","pg_dump → RDS\nreset de secuencias",       RED),
        ("RDS · Cloud", "PostgreSQL en AWS\nTablas META + FACT",    GREEN),
        ("Pipeline",    "Métricas · LLM · PDF\nS3 + SES email",     ACCENT),
    ]

    node_w = Inches(2.65)
    node_h = Inches(1.55)
    gap    = Inches(0.55)
    total  = len(nodes) * node_w + (len(nodes) - 1) * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.1)

    for i, (title, desc, color) in enumerate(nodes):
        x = start_x + i * (node_w + gap)
        _rect(s, x, y, node_w, node_h, BG_CARD, line_color=color, line_pt=1)
        _rect(s, x, y, node_w, Inches(0.04), color)
        _txt(s, title, x + Inches(0.14), y + Inches(0.12),
             node_w - Inches(0.28), Inches(0.5),
             size=12, bold=True, color=color)
        _txt(s, desc, x + Inches(0.14), y + Inches(0.65),
             node_w - Inches(0.28), Inches(0.82),
             size=9.5, color=TEXT_GRAY)

        if i < len(nodes) - 1:
            ax = x + node_w + Inches(0.1)
            _txt(s, "→", ax, y + Inches(0.55),
                 gap - Inches(0.2), Inches(0.4),
                 size=16, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Alerta de pain point
    _rect(s, Inches(0.5), Inches(3.9), SLIDE_W - Inches(1.0), Inches(0.52),
          RED_DARK, line_color=RED, line_pt=0.75)
    _txt(s, "⚠  El paso SYNC MANUAL es el punto de falla central — sin él el pipeline no corre y el monitoreo se interrumpe sin aviso.",
         Inches(0.68), Inches(3.95), SLIDE_W - Inches(1.35), Inches(0.42),
         size=10.5, color=RED_SOFT)

    _bullets(s, [
        ("Coordinación semanal obligatoria",
         "el equipo de Crédito debe ejecutar el sync antes de que el pipeline pueda correr"),
        ("Reset manual de secuencias SQL",
         "tras cada pg_dump se requiere SELECT setval() por tabla o los inserts del pipeline rompen"),
        ("Acoplamiento de versiones",
         "ETL y Pipeline deben ser compatibles en schema; un deploy asimétrico genera errores silenciosos"),
    ], Inches(0.5), Inches(4.6), SLIDE_W - Inches(1.0), spacing=0.6, size=11)

    _page_num(s, 2)


# ---------------------------------------------------------------------------
# Slide 3 — Dependencia operativa bilateral
# ---------------------------------------------------------------------------

def slide_dependencia(prs):
    s = _blank_slide(prs)
    _badge(s, "PROBLEMA 1 · DEPENDENCIA OPERATIVA")

    _txt(s, "El pipeline depende de un actor externo",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    _txt(s, "Cualquier semana en que el sync falla es una semana sin monitoreo",
         Inches(0.5), Inches(1.42), Inches(12), Inches(0.4),
         size=12, color=TEXT_DIM)

    _divider(s)

    _bullets(s, [
        ("Sync manual semanal",
         "el pipeline en nube no corre hasta que Crédito ejecute pg_dump/pg_restore a RDS"),
        ("Sin señal automática de fallo",
         "si el sync se omite, el reporte se genera con datos viejos sin ninguna advertencia visible"),
        ("Cambios de schema = doble coordinación",
         "nuevo segmento o target requiere re-ejecutar bootstrap en CML y volcado completo a RDS"),
        ("Reset de secuencias SQL frágil",
         "tras cada volcado hay que correr SELECT setval() por tabla o los inserts del Pipeline rompen por colisión de IDs"),
        ("Acoplamiento de versiones entre entornos",
         "si ETL (CML) y Pipeline (nube) se actualizan por separado, las incompatibilidades de schema son silenciosas y difíciles de depurar"),
    ], Inches(0.5), Inches(2.08), SLIDE_W - Inches(1.0), spacing=0.9, size=12)

    _page_num(s, 3)


# ---------------------------------------------------------------------------
# Slide 4 — Complejidad técnica: CI/CD + Impala (dos columnas)
# ---------------------------------------------------------------------------

def slide_complejidad(prs):
    s = _blank_slide(prs)
    _badge(s, "PROBLEMA 2 · COMPLEJIDAD TÉCNICA")

    _txt(s, "Dos frentes de complejidad simultáneos",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    _txt(s, "CI/CD roto  ·  motor de BD incompatible con el stack actual",
         Inches(0.5), Inches(1.42), Inches(12), Inches(0.4),
         size=12, color=TEXT_DIM)

    _divider(s)

    col_w   = Inches(5.95)
    left1   = Inches(0.5)
    left2   = Inches(7.0)
    top_hdr = Inches(2.05)
    top_bul = Inches(2.55)

    # Column headers
    _txt(s, "Desarrollo y CI/CD", left1, top_hdr, col_w, Inches(0.38),
         size=11, bold=True, color=ACCENT_LIGHT)
    _txt(s, "Impala como motor de datos en CML", left2, top_hdr, col_w, Inches(0.38),
         size=11, bold=True, color=YELLOW)

    # Separador vertical
    _rect(s, Inches(6.65), Inches(2.0), Inches(0.018), Inches(4.8), TEXT_DIM)

    _bullets(s, [
        ("CML sin acceso a GitHub",
         "código se actualiza por SFTP/VPN; sin CI/CD ni trazabilidad automática"),
        ("Code drift inevitable",
         "sin automatización, ETL y Pipeline divergen entre deploys sin que nadie lo note"),
        ("Entorno de prueba no replicable",
         "no hay forma práctica de simular Impala en local para testear el ETL antes de subir a CML"),
        ("Dos perfiles de dependencias",
         "cada cambio se valida en dos entornos con stacks distintos, duplicando el esfuerzo de QA"),
    ], left1, top_bul, col_w, spacing=0.98, size=11, accent_color=ACCENT_LIGHT)

    _bullets(s, [
        ("SQLAlchemy no soporta Impala",
         "no hay dialect oficial; las alternativas de terceros están abandonadas"),
        ("Sin ACID completo ni sequences",
         "idempotencia del ETL y llaves primarias auto-increment requieren reescritura completa"),
        ("Sin migraciones con Alembic",
         "cambios de schema son manuales, sin trazabilidad ni posibilidad de rollback"),
        ("Dev/prod parity desaparece",
         "hoy SQLite local = PostgreSQL prod con el mismo models.py; con Impala ese contrato se rompe"),
    ], left2, top_bul, col_w, spacing=0.98, size=11, accent_color=YELLOW)

    _page_num(s, 4)


# ---------------------------------------------------------------------------
# Slide 5 — Escalar en este esquema es insostenible
# ---------------------------------------------------------------------------

def slide_escala(prs):
    s = _blank_slide(prs)
    _badge(s, "PROBLEMA 3 · ESCALABILIDAD")

    _txt(s, "Añadir modelos crece en complejidad de forma no lineal",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    _txt(s, "Cada nuevo modelo o equipo multiplica las dependencias manuales",
         Inches(0.5), Inches(1.42), Inches(12), Inches(0.4),
         size=12, color=TEXT_DIM)

    _divider(s)

    _bullets(s, [
        ("Colisión de IDs entre equipos",
         "bootstrap independiente por equipo genera llaves primarias duplicadas en la misma RDS"),
        ("Sin gobernanza del catálogo de modelos",
         "meta_model_registry gestionado por cada equipo en su CML = inconsistencia directa en dashboards BI"),
        ("Onboarding lineal en trabajo manual",
         "cada nuevo modelo suma bootstrap + ETL + sync + reset de secuencias al ciclo semanal"),
        ("Dependencias cruzadas insostenibles",
         "el pipeline en nube queda bloqueado por N equipos con sync independiente; un fallo paraliza todo"),
        ("Sin validación de integridad del sync",
         "no hay garantía automática de que los datos llegaron completos y consistentes a RDS"),
    ], Inches(0.5), Inches(2.08), SLIDE_W - Inches(1.0), spacing=0.9, size=12)

    _page_num(s, 5)


# ---------------------------------------------------------------------------
# Slide 6 — Trabajo con fecha de caducidad
# ---------------------------------------------------------------------------

def slide_caducidad(prs):
    s = _blank_slide(prs)
    _badge(s, "PROBLEMA 4 · COSTO DE OPORTUNIDAD")

    _txt(s, "Todo este trabajo tiene fecha de caducidad",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    _txt(s, "La empresa migra todo a AWS en meses — CML será abandonado",
         Inches(0.5), Inches(1.42), Inches(12), Inches(0.4),
         size=12, color=TEXT_DIM)

    _divider(s)

    _bullets(s, [
        ("La adaptación para CML quedará obsoleta",
         "en el momento en que se abandone CML, todo el trabajo de separar ETL/Pipeline se vuelve innecesario"),
        ("La herramienta regresa casi a su estado actual",
         "sin CML, el flujo vuelve a ser: datos en nube → ETL → BD → Pipeline → Reporte; sin intermediarios"),
        ("El sync manual desaparece en 100% cloud",
         "el ETL puede leer directamente de la fuente de datos sin pasos manuales del equipo de Crédito"),
        ("Costo de oportunidad real",
         "tiempo invertido en adaptar a CML = tiempo que no se dedica a conectar datos reales, escalar modelos o construir BI"),
    ], Inches(0.5), Inches(2.08), SLIDE_W - Inches(1.0), spacing=0.9, size=12)

    # Caja de cierre
    box_y = Inches(5.85)
    _rect(s, Inches(0.5), box_y, SLIDE_W - Inches(1.0), Inches(0.78),
          RGBColor(0x1C, 0x14, 0x06), line_color=ACCENT, line_pt=1)
    _txt(s, "Recomendación: desplegar directamente en AWS y evitar el esquema híbrido.",
         Inches(0.72), box_y + Inches(0.06), SLIDE_W - Inches(1.44), Inches(0.36),
         size=13, bold=True, color=ACCENT_LIGHT)
    _txt(s, "El stack ya está listo — solo se necesita la fuente de datos en nube.",
         Inches(0.72), box_y + Inches(0.4), SLIDE_W - Inches(1.44), Inches(0.3),
         size=11, color=TEXT_GRAY)

    _page_num(s, 6)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def create_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_portada(prs)
    slide_arquitectura(prs)
    slide_dependencia(prs)
    slide_complejidad(prs)
    slide_escala(prs)
    slide_caducidad(prs)

    output_path = Path(__file__).parent.parent / "artifacts" / "ppt" / "MLMonitor_CML_Decision.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"✓  Presentación guardada en: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
