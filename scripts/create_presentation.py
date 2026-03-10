"""
Script para generar la presentación PowerPoint "ML Monitor Weekly".
Ejecutar con: poetry run python scripts/create_presentation.py
Output: artifacts/mlmonitor_weekly_presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paleta estilo Anthropic
# ---------------------------------------------------------------------------
BG_DARK      = RGBColor(0x0F, 0x0F, 0x0E)
BG_CARD      = RGBColor(0x1C, 0x1C, 0x1A)
ACCENT       = RGBColor(0xD9, 0x77, 0x06)
ACCENT_LIGHT = RGBColor(0xF5, 0xA6, 0x23)
TEXT_CREAM   = RGBColor(0xE8, 0xE3, 0xDB)
TEXT_GRAY    = RGBColor(0xB0, 0xAD, 0xA8)
TEXT_DIM     = RGBColor(0x6B, 0x68, 0x65)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
YELLOW       = RGBColor(0xF5, 0x9E, 0x0B)
RED          = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide


def _txt(slide, text, left, top, width, height,
         size=14, bold=False, color=TEXT_CREAM,
         align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box


def _rect(slide, left, top, width, height, fill_color,
          line_color=None, line_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _code(slide, code_text, left, top, width, height):
    """Bloque de código estilo terminal con borde naranja."""
    _rect(slide, left, top, width, height, BG_CARD,
          line_color=ACCENT, line_pt=0.75)

    box = slide.shapes.add_textbox(
        left + Inches(0.14), top + Inches(0.1),
        width - Inches(0.28), height - Inches(0.2),
    )
    tf = box.text_frame
    tf.word_wrap = False

    for i, line in enumerate(code_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(8.5)
        r.font.name = "Courier New"
        if line.strip().startswith("#"):
            r.font.color.rgb = TEXT_DIM
        elif any(kw in line for kw in ("def ", "class ", "import ", "from ")):
            r.font.color.rgb = ACCENT_LIGHT
        elif "=" in line and "==" not in line:
            r.font.color.rgb = TEXT_CREAM
        else:
            r.font.color.rgb = TEXT_GRAY


def _badge(slide, label):
    """Pequeño badge de sección en esquina superior izquierda."""
    _rect(slide, Inches(0.5), Inches(0.28), Inches(2.8), Inches(0.32),
          RGBColor(0x2A, 0x1E, 0x08), line_color=ACCENT, line_pt=0.5)
    _txt(slide, label,
         Inches(0.58), Inches(0.28), Inches(2.7), Inches(0.32),
         size=7.5, bold=True, color=ACCENT)


def _page_num(slide, n):
    _txt(slide, str(n),
         SLIDE_W - Inches(0.6), SLIDE_H - Inches(0.38),
         Inches(0.5), Inches(0.3),
         size=9, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 1 — Título
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = _blank_slide(prs)

    # Barra izquierda naranja
    _rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)

    _txt(s, "ML Monitor Weekly",
         Inches(1.2), Inches(1.9), Inches(11), Inches(1.4),
         size=52, bold=True, color=TEXT_CREAM)

    _txt(s, "Avances de integración con AWS",
         Inches(1.2), Inches(3.25), Inches(11), Inches(0.7),
         size=22, color=ACCENT)

    # Línea separadora
    _rect(s, Inches(1.2), Inches(3.95), Inches(4.5), Inches(0.02), TEXT_DIM)

    _txt(s, "Marzo 2026  ·  BazBoost Credit Scorecard Fleet",
         Inches(1.2), Inches(4.1), Inches(11), Inches(0.5),
         size=13, color=TEXT_DIM)

    # Pills de servicios AWS
    services = ["Amazon RDS", "Amazon Bedrock", "Amazon S3", "Amazon SES"]
    for i, svc in enumerate(services):
        x = Inches(1.2 + i * 2.9)
        _rect(s, x, Inches(5.35), Inches(2.6), Inches(0.48),
              BG_CARD, line_color=ACCENT, line_pt=0.75)
        _txt(s, svc, x, Inches(5.35), Inches(2.6), Inches(0.48),
             size=11, bold=True, color=TEXT_CREAM, align=PP_ALIGN.CENTER)

    _page_num(s, 1)


# ---------------------------------------------------------------------------
# Slide 2 — Arquitectura del Pipeline
# ---------------------------------------------------------------------------

def slide_architecture(prs):
    s = _blank_slide(prs)
    _badge(s, "ARQUITECTURA")

    _txt(s, "Pipeline de Monitoreo",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    _txt(s, "4 pasos · 4 servicios AWS · ejecución semanal automatizada",
         Inches(0.5), Inches(1.42), Inches(12), Inches(0.4),
         size=12, color=TEXT_DIM)

    steps = [
        ("01", "Calcular\nMétricas",    "PSI · Gini · KS\nNull rates · RollForward", "Amazon RDS",     GREEN),
        ("02", "Analizar\ncon LLM",     "Narrativas por segmento\nAcciones recomendadas", "Amazon Bedrock", ACCENT),
        ("03", "Generar\nPDF",          "Jinja2 → WeasyPrint\nHTML → PDF (A4)",       "Artifact local", ACCENT_LIGHT),
        ("04", "Distribuir\nReporte",   "Upload S3 + Envío email\ncon adjunto PDF",   "S3 + SES",       YELLOW),
    ]

    for i, (num, title, desc, svc, color) in enumerate(steps):
        x = Inches(0.35 + i * 3.2)
        y = Inches(2.05)
        w = Inches(2.95)
        h = Inches(4.4)

        _rect(s, x, y, w, h, BG_CARD, line_color=color, line_pt=1)

        _txt(s, num,   x + Inches(0.15), y + Inches(0.15), Inches(0.7), Inches(0.5),
             size=22, bold=True, color=color)
        _txt(s, title, x + Inches(0.15), y + Inches(0.65), Inches(2.6), Inches(0.9),
             size=16, bold=True)
        _txt(s, desc,  x + Inches(0.15), y + Inches(1.6),  Inches(2.6), Inches(1.4),
             size=11, color=TEXT_GRAY)

        # Badge de servicio en fondo del card
        _rect(s, x + Inches(0.1), y + h - Inches(0.55), w - Inches(0.2), Inches(0.4),
              RGBColor(0x0A, 0x0A, 0x09))
        _txt(s, svc,
             x + Inches(0.1), y + h - Inches(0.56), w - Inches(0.2), Inches(0.4),
             size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Flecha entre cards
        if i < 3:
            _txt(s, "→",
                 x + w + Inches(0.04), y + h / 2 - Inches(0.18),
                 Inches(0.22), Inches(0.36),
                 size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    _page_num(s, 2)


# ---------------------------------------------------------------------------
# Slide 3 — RDS: Modelo de datos
# ---------------------------------------------------------------------------

def slide_rds_model(prs):
    s = _blank_slide(prs)
    _badge(s, "AMAZON RDS · POSTGRESQL")

    _txt(s, "Modelo de Datos",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    _txt(s, "6 tablas · Patrón SCD2 para catálogos · Append-only para métricas históricas",
         Inches(0.5), Inches(1.42), Inches(12), Inches(0.4),
         size=12, color=TEXT_DIM)

    groups = [
        ("META — Catálogos (SCD2)", ACCENT_LIGHT, [
            ("meta_model_registry",
             "Catálogo maestro de modelos\nmodel_id · fleet_id · score_min/max · lag_semanas · valid_from/valid_to"),
            ("meta_variables",
             "Catálogo de variables por modelo\nvariable_name · type (numeric|categorical) · binning_rules (JSON)"),
            ("meta_metric_thresholds",
             "Umbrales de alerta por métrica\nwarning_threshold · critical_threshold · direction"),
        ]),
        ("FACT — Historial (Append-only)", GREEN, [
            ("fact_distributions",
             "Distribuciones de variables por semana\nbin_label · bin_count · null_count · bin_percentage"),
            ("fact_performance_outcomes",
             "Outcomes de performance por decil\nroll_forward · payment_rate · calibración (Gini/KS)"),
            ("fact_metrics_history",
             "Historial de métricas calculadas\nmetric_value · alert_label: OK | WARNING | CRITICAL"),
        ]),
    ]

    for col, (group_name, color, tables) in enumerate(groups):
        x = Inches(0.5 + col * 6.4)

        _txt(s, group_name, x, Inches(1.98), Inches(6.1), Inches(0.38),
             size=11, bold=True, color=color)

        for i, (tname, desc) in enumerate(tables):
            ty = Inches(2.42 + i * 1.6)
            cw = Inches(5.95)
            ch = Inches(1.45)

            _rect(s, x, ty, cw, ch, BG_CARD, line_color=color, line_pt=0.75)
            # acento izquierdo
            _rect(s, x, ty, Inches(0.04), ch, color)

            _txt(s, tname,
                 x + Inches(0.18), ty + Inches(0.1), Inches(5.5), Inches(0.38),
                 size=10.5, bold=True, color=TEXT_CREAM, font="Courier New")
            _txt(s, desc,
                 x + Inches(0.18), ty + Inches(0.52), Inches(5.5), Inches(0.8),
                 size=9.5, color=TEXT_GRAY)

    _page_num(s, 3)


# ---------------------------------------------------------------------------
# Slide 4 — RDS: Consultas
# ---------------------------------------------------------------------------

def slide_rds_queries(prs):
    s = _blank_slide(prs)
    _badge(s, "AMAZON RDS · CONSULTAS EN VIVO")

    _txt(s, "Datos en RDS PostgreSQL",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    # --- Query 1 ---
    _txt(s, "▸  Catálogo de modelos activos (meta_model_registry)",
         Inches(0.5), Inches(1.58), Inches(12), Inches(0.35),
         size=10.5, bold=True, color=ACCENT_LIGHT)

    q1 = ("SELECT model_id, fleet_id, model_type,\n"
          "       score_min, score_max, lag_semanas\n"
          "FROM meta_model_registry\n"
          "WHERE valid_to IS NULL;")
    _code(s, q1, Inches(0.5), Inches(1.95), Inches(5.7), Inches(1.05))

    # Placeholder resultado query 1
    _rect(s, Inches(6.4), Inches(1.58), Inches(6.5), Inches(1.42),
          RGBColor(0x1A, 0x1A, 0x18), line_color=TEXT_DIM, line_pt=0.5)
    _txt(s, "[ Pegar resultado de la consulta aquí ]",
         Inches(6.55), Inches(1.75), Inches(6.2), Inches(0.9),
         size=10, color=TEXT_DIM, font="Courier New")

    # --- Query 2 ---
    _txt(s, "▸  Alertas activas · semana 2026-01-05 (fact_metrics_history)",
         Inches(0.5), Inches(3.2), Inches(12), Inches(0.35),
         size=10.5, bold=True, color=ACCENT_LIGHT)

    q2 = ("SELECT segment_id, metric_name,\n"
          "       ROUND(metric_value::numeric, 4) AS valor,\n"
          "       alert_label, calculation_week\n"
          "FROM fact_metrics_history\n"
          "WHERE calculation_week = '2026-01-05'\n"
          "  AND alert_label IN ('CRITICAL','WARNING')\n"
          "ORDER BY alert_label, metric_name\n"
          "LIMIT 15;")
    _code(s, q2, Inches(0.5), Inches(3.57), Inches(5.7), Inches(1.95))

    # Placeholder resultado query 2
    _rect(s, Inches(6.4), Inches(3.2), Inches(6.5), Inches(2.32),
          RGBColor(0x1A, 0x1A, 0x18), line_color=TEXT_DIM, line_pt=0.5)
    _txt(s, "[ Pegar resultado de la consulta aquí ]",
         Inches(6.55), Inches(3.9), Inches(6.2), Inches(1.2),
         size=10, color=TEXT_DIM, font="Courier New")

    # Nota al pie
    _txt(s, "✦  18,035 filas en producción  ·  5 tablas  ·  Semana de referencia: 2026-01-05  ·  DB: RDS PostgreSQL",
         Inches(0.5), Inches(5.7), Inches(12.5), Inches(0.4),
         size=9.5, color=TEXT_DIM)

    _page_num(s, 4)


# ---------------------------------------------------------------------------
# Slide 5 — Bedrock
# ---------------------------------------------------------------------------

def slide_bedrock(prs):
    s = _blank_slide(prs)
    _badge(s, "AMAZON BEDROCK · LLM")

    _txt(s, "Análisis Automático con LLM",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    # Badge de modelo
    _rect(s, Inches(0.5), Inches(1.55), Inches(8.2), Inches(0.42),
          RGBColor(0x1C, 0x14, 0x06), line_color=ACCENT, line_pt=0.75)
    _txt(s, "modelo:  us.anthropic.claude-haiku-4-5-20251001-v1:0  ·  bedrock-runtime  ·  us-east-1",
         Inches(0.65), Inches(1.57), Inches(8.0), Inches(0.38),
         size=9, color=ACCENT, font="Courier New")

    # Snippet de código
    snippet = ("# bedrock_analyst.py\n"
               "def _call_llm(self, prompt: str) -> str:\n"
               "    client = self._get_client()\n"
               "    body = json.dumps({\n"
               '        "anthropic_version": "bedrock-2023-05-31",\n'
               '        "max_tokens": 2048,\n'
               '        "temperature": 0.2,\n'
               '        "messages": [{"role": "user", "content": prompt}],\n'
               "    })\n"
               "    response = client.invoke_model(\n"
               "        modelId=self.model_id,\n"
               "        body=body,\n"
               '        contentType="application/json",\n'
               "    )\n"
               '    result = json.loads(response["body"].read())\n'
               '    return result["content"][0]["text"]')
    _code(s, snippet, Inches(0.5), Inches(2.1), Inches(6.5), Inches(4.55))

    # ¿Qué genera?
    _txt(s, "¿Qué genera el LLM?",
         Inches(7.2), Inches(2.1), Inches(5.8), Inches(0.4),
         size=13, bold=True)

    outputs = [
        (ACCENT,       "Fleet Narrative",
         "1 llamada a Bedrock\nAnálisis ejecutivo del estado de la flota completa"),
        (GREEN,        "Segment Narratives",
         "9 llamadas individuales\nDiagnóstico por segmento con contexto de métricas PSI/Gini/KS"),
        (YELLOW,       "Recommended Actions",
         "JSON estructurado por segmento\nPrioridad: CRÍTICO · ALTA · MEDIA · MONITOREO"),
    ]

    for i, (color, title, desc) in enumerate(outputs):
        oy = Inches(2.65 + i * 1.35)
        _rect(s, Inches(7.2), oy, Inches(5.8), Inches(1.2), BG_CARD,
              line_color=color, line_pt=0.75)
        _rect(s, Inches(7.2), oy, Inches(0.04), Inches(1.2), color)
        _txt(s, title,
             Inches(7.38), oy + Inches(0.1), Inches(5.3), Inches(0.35),
             size=11, bold=True, color=color)
        _txt(s, desc,
             Inches(7.38), oy + Inches(0.48), Inches(5.3), Inches(0.62),
             size=9.5, color=TEXT_GRAY)

    _page_num(s, 5)


# ---------------------------------------------------------------------------
# Slide 6 — S3
# ---------------------------------------------------------------------------

def slide_s3(prs):
    s = _blank_slide(prs)
    _badge(s, "AMAZON S3 · ALMACENAMIENTO")

    _txt(s, "Reportes PDF en S3",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    # URI
    _txt(s, "URI del reporte en producción:",
         Inches(0.5), Inches(1.58), Inches(5), Inches(0.35),
         size=10, color=TEXT_DIM)

    _rect(s, Inches(0.5), Inches(1.93), Inches(12.3), Inches(0.5),
          BG_CARD, line_color=ACCENT, line_pt=1)
    _txt(s, "s3://ml-monitoring-reports-credito/mlmonitor/reports/mlmonitor_2026-01-05.pdf",
         Inches(0.65), Inches(1.96), Inches(12.0), Inches(0.42),
         size=10.5, color=ACCENT, font="Courier New")

    # Código S3Uploader
    s3_code = ("# s3_uploader.py\n"
               "def upload(self, local_path: Path) -> str:\n"
               "    key = f\"{self.prefix}/{local_path.name}\"\n"
               "    self._get_client().upload_file(\n"
               "        Filename=str(local_path),\n"
               "        Bucket=self.bucket,\n"
               "        Key=key,\n"
               '        ExtraArgs={"ContentType": "application/pdf"},\n'
               "    )\n"
               "    s3_uri = f\"s3://{self.bucket}/{key}\"\n"
               "    print(f\"[S3Uploader] Subido: {s3_uri}\")\n"
               "    return s3_uri")
    _code(s, s3_code, Inches(0.5), Inches(2.58), Inches(5.95), Inches(3.65))

    # Output real de aws s3 ls
    _txt(s, "$ aws s3 ls s3://ml-monitoring-reports-credito/mlmonitor/reports/",
         Inches(6.65), Inches(2.58), Inches(6.4), Inches(0.38),
         size=9, color=TEXT_DIM, font="Courier New")

    _rect(s, Inches(6.65), Inches(2.96), Inches(6.3), Inches(0.58),
          BG_CARD, line_color=GREEN, line_pt=0.75)
    _txt(s, "2026-03-04 12:35:24     117899 mlmonitor_2026-01-05.pdf",
         Inches(6.8), Inches(3.01), Inches(6.0), Inches(0.45),
         size=9.5, color=GREEN, font="Courier New")

    # Placeholder captura pantalla
    _rect(s, Inches(6.65), Inches(3.68), Inches(6.3), Inches(2.55),
          RGBColor(0x14, 0x14, 0x12), line_color=TEXT_DIM, line_pt=0.5)
    _txt(s, "[ Insertar captura de pantalla\ndel objeto en S3 Console\no URL del PDF abierto ]",
         Inches(7.0), Inches(4.4), Inches(5.6), Inches(1.3),
         size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    _page_num(s, 6)


# ---------------------------------------------------------------------------
# Slide 7 — SES
# ---------------------------------------------------------------------------

def slide_ses(prs):
    s = _blank_slide(prs)
    _badge(s, "AMAZON SES · EMAIL")

    _txt(s, "Distribución Automática por Email",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=30, bold=True)

    # Código SES
    ses_code = ("# sender.py — SESEmailSender\n"
                "self._get_client().send_raw_email(\n"
                "    Source=self.from_email,\n"
                "    Destinations=recipients,\n"
                "    RawMessage={\"Data\": msg.as_string()},\n"
                ")\n"
                "\n"
                "# msg = MIMEMultipart con:\n"
                "#   MIMEText        → cuerpo HTML branding MLMonitor\n"
                "#   MIMEApplication → adjunto PDF\n"
                "#     (mlmonitor_2026-01-05.pdf, 115 KB)")
    _code(s, ses_code, Inches(0.5), Inches(1.6), Inches(5.8), Inches(3.1))

    # Anatomía del correo
    _txt(s, "Anatomía del correo enviado:",
         Inches(0.5), Inches(4.85), Inches(6.0), Inches(0.35),
         size=10, bold=True, color=TEXT_DIM)

    fields = [
        ("From:",    "MLMonitor <noreply@mlmonitor.com>"),
        ("To:",      "Lista configurada en SES_RECIPIENTS (env var)"),
        ("Subject:", "MLMonitor — Reporte de Monitoreo mlmonitor_2026-01-05"),
        ("Adjunto:", "mlmonitor_2026-01-05.pdf  ·  115 KB  ·  application/pdf"),
    ]

    for i, (label, value) in enumerate(fields):
        fy = Inches(5.25 + i * 0.38)
        _txt(s, label, Inches(0.5), fy, Inches(1.1), Inches(0.35),
             size=9, bold=True, color=ACCENT_LIGHT, font="Courier New")
        _txt(s, value, Inches(1.65), fy, Inches(5.2), Inches(0.35),
             size=9, color=TEXT_GRAY, font="Courier New")

    # Placeholder captura email
    _rect(s, Inches(6.5), Inches(1.55), Inches(6.5), Inches(5.6),
          RGBColor(0x14, 0x14, 0x12), line_color=TEXT_DIM, line_pt=0.5)
    _txt(s, "[ Insertar captura de pantalla\ndel correo recibido\nen la bandeja de entrada ]",
         Inches(7.0), Inches(3.7), Inches(5.5), Inches(1.5),
         size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    _page_num(s, 7)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def create_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_architecture(prs)
    slide_rds_model(prs)
    slide_rds_queries(prs)
    slide_bedrock(prs)
    slide_s3(prs)
    slide_ses(prs)

    output_dir = Path(__file__).parent.parent / "artifacts"
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "mlmonitor_weekly_presentation.pptx"
    prs.save(str(output_path))
    print(f"✓  Presentación guardada en: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
